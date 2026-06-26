"""
Slide generation for Pixel Assistant.
Requires: pip install python-pptx

Usage (internal): generate_slides(title, bullet_sections, theme, out_dir) -> Path
"""
import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt


# ── Themes ─────────────────────────────────────────────────────────────────────

THEMES = {
    "dark": {
        "bg":        RGBColor(0x0F, 0x14, 0x28),
        "title_fg":  RGBColor(0x40, 0x80, 0xFF),
        "body_fg":   RGBColor(0xDC, 0xEB, 0xFF),
        "accent":    RGBColor(0x40, 0x80, 0xFF),
    },
    "light": {
        "bg":        RGBColor(0xFF, 0xFF, 0xFF),
        "title_fg":  RGBColor(0x1A, 0x1A, 0x2E),
        "body_fg":   RGBColor(0x33, 0x33, 0x33),
        "accent":    RGBColor(0x28, 0x78, 0xC8),
    },
    "corporate": {
        "bg":        RGBColor(0xF5, 0xF7, 0xFA),
        "title_fg":  RGBColor(0x00, 0x33, 0x66),
        "body_fg":   RGBColor(0x22, 0x22, 0x22),
        "accent":    RGBColor(0x00, 0x66, 0xCC),
    },
    "modern": {
        "bg":        RGBColor(0x1E, 0x20, 0x30),
        "title_fg":  RGBColor(0x00, 0xE5, 0xFF),
        "body_fg":   RGBColor(0xE0, 0xE0, 0xE0),
        "accent":    RGBColor(0x00, 0xE5, 0xFF),
    },
    "warm": {
        "bg":        RGBColor(0xFF, 0xF8, 0xF0),
        "title_fg":  RGBColor(0x8B, 0x45, 0x13),
        "body_fg":   RGBColor(0x3B, 0x1A, 0x00),
        "accent":    RGBColor(0xD2, 0x69, 0x1E),
    },
}


def _apply_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_text(tf, text: str, size: int, color: RGBColor, bold: bool = False):
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size  = Pt(size)
            run.font.color.rgb = color
            run.font.bold  = bold


def generate_slides(
    title: str,
    sections: list[dict],   # [{"heading": str, "bullets": [str]}]
    theme: str = "dark",
    out_dir: Path = None,
) -> Path:
    """
    sections = [{"heading": "Intro", "bullets": ["Point 1", "Point 2"]}, ...]
    Returns the path to the saved .pptx file.
    """
    t = THEMES.get(theme, THEMES["dark"])
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    # ── Title slide ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _apply_bg(slide, t["bg"])

    tf = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
    _set_text(tf.text_frame, title, 44, t["title_fg"], bold=True)

    # accent line
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(1), Inches(4.2), Inches(3), Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = t["accent"]
    line.line.fill.background()

    # ── Content slides ────────────────────────────────────────────────────
    for section in sections:
        slide = prs.slides.add_slide(blank_layout)
        _apply_bg(slide, t["bg"])

        # Heading
        heading_tf = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.9)
        )
        _set_text(heading_tf.text_frame, section.get("heading", ""), 32, t["title_fg"], bold=True)

        # Accent bar under heading
        bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.25), Inches(12), Inches(0.04))
        bar.fill.solid()
        bar.fill.fore_color.rgb = t["accent"]
        bar.line.fill.background()

        # Bullets
        body_tf = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5)
        )
        tf_frame = body_tf.text_frame
        tf_frame.word_wrap = True

        bullets = section.get("bullets", [])
        for i, bullet in enumerate(bullets):
            if i == 0:
                para = tf_frame.paragraphs[0]
            else:
                para = tf_frame.add_paragraph()
            para.text = f"• {bullet}"
            para.space_before = Pt(6)
            for run in para.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = t["body_fg"]

    if out_dir is None:
        out_dir = Path(__file__).parent.parent.parent / "generated"
    out_dir.mkdir(exist_ok=True)

    safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:40]
    out_path = out_dir / f"{safe_title}_{theme}.pptx"
    prs.save(str(out_path))
    return out_path
